#!/usr/bin/env python3
"""Build the ENPM690 final report package.

Outputs report Markdown, PDF, PPTX, demo scripts, validation script, and
summary artifacts. All headline values are centralized here and mirrored in
report/OFFICIAL_ARTIFACTS.md.
"""

from __future__ import annotations

import json
import os
import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import Image, PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle


ROOT = Path(__file__).resolve().parents[2]
REPORT = ROOT / "report"
FIG = REPORT / "figures"
SCRIPTS = ROOT / "scripts" / "final"


OFFICIAL = {
    "title": "Modular L1-to-RL Manipulation Stack with Route-Curriculum Extension",
    "stage_rows": [
        [0, 1.00, 0.50, 0.0073, 1.67, 0.0106],
        [1, 1.00, 0.62, 0.0099, 1.67, 0.0123],
        [2, 1.00, 0.85, 0.0119, 1.82, 0.0139],
        [3, 1.00, 1.20, 0.0138, 2.14, 0.0164],
        [4, 1.00, 1.71, 0.0150, 2.53, 0.0165],
        [5, 0.93, 1.96, 0.0177, 2.89, 0.0208],
    ],
    "qwen_command": "Move tray1 from shelf_A1 to shelf_B1 while keeping it level and inserting with a stable pose.",
    "route": {
        "baseline_success": 0.0435,
        "baseline_prefix": 21,
        "best_checkpoint": "artifacts/kinematic_phase1/route_curriculum/route_prefix120_routeobs_sequence2_1m_001/model_latest.zip",
        "prefix120_success": 1.0,
        "prefix120_prefix": 120,
        "prefix120_distance": 1.720,
        "prefix120_pos_error": 0.00934,
        "prefix120_ori_error": 0.02444,
        "full483_success": 0.4741,
        "full483_prefix": 170,
        "full483_distance": 2.455,
        "first_failure_index": 171,
        "first_failure_reason": "position",
    },
}


def read_json(path: str):
    p = ROOT / path
    if not p.exists():
        return None
    try:
        return json.loads(p.read_text())
    except Exception:
        return None


def ensure_dirs() -> None:
    for path in [REPORT, FIG, REPORT / "demo_outputs", REPORT / "videos", SCRIPTS]:
        path.mkdir(parents=True, exist_ok=True)


def write(path: Path, text: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(text.strip() + "\n", encoding="utf-8")


def exists(path: str) -> bool:
    return (ROOT / path).exists()


def audit_md() -> str:
    report_files = sorted(str(p.relative_to(ROOT)) for p in (ROOT / "docs").glob("*.md"))
    official_paths = [
        "docs/PHASE2_FINAL_DEMONSTRATION_REPORT.md",
        "docs/RL_WORKSPACE_AND_TRANSPORT_STATUS.md",
        "docs/ROUTE_CURRICULUM_TRAINING_PLAN.md",
        "artifacts/kinematic_phase1/phase1c/workspace_sweep_workspace_noop_vs_previous_summary_001.json",
        "artifacts/v5/qwen_l1_demo/l1_to_rl_skill_request_qwen.json",
        OFFICIAL["route"]["best_checkpoint"],
        "artifacts/kinematic_phase1/route_curriculum/prefix120_teacher_anchor/teacher_route_anchor_dataset.npz",
        "artifacts/kinematic_phase1/route_curriculum/route_segment121_180_teacheranchored_smoke_001/route_gate_full483/route_gate_summary.json",
    ]
    found = [p for p in official_paths if exists(p)]
    missing = [p for p in official_paths if not exists(p)]
    return f"""
# Final Package Audit

Purpose: Audit the current repository state before building the final ENPM690 report package.

## Existing Report / Documentation Files

The current repo primarily stores research documents in `docs/`; the final package is generated into `report/`.

{chr(10).join(f'- `{p}`' for p in report_files[:80])}

## Official Artifacts Found

{chr(10).join(f'- `{p}`' for p in found)}

## Missing Or Optional Artifacts

{chr(10).join(f'- `{p}`' for p in missing) if missing else '- None for required final-package artifacts.'}

## Files To Regenerate For Final Package

- `report/FINAL_REPORT.md`
- `report/FINAL_REPORT.pdf`
- `report/FINAL_PRESENTATION.pptx`
- `report/figures/*.png`
- `report/demo_outputs/*.json`

## Artifacts Not To Use As Official Best

- `artifacts/kinematic_phase1/route_curriculum/route_prefix180_routeobs_sequence2_1m_001/model_latest.zip`
- `artifacts/kinematic_phase1/route_curriculum/route_prefix180_routeobs_sequence2_antiforget_1m_001/model_latest.zip`
- `artifacts/kinematic_phase1/route_curriculum/route_segment121_180_teacheranchored_1m_001/model_latest.zip`

These are useful negative-result artifacts, but they fail sequential retention and should not replace the official prefix120 checkpoint.

## Official Numbers For Final Package

- Stage 5 kinematic skill success: `0.93`
- Stage 5 final position error: `2.89 mm`
- Stage 5 final orientation error: `0.0208 rad`
- Route baseline full483 success: `0.0435`
- Route baseline longest prefix: `21`
- Route prefix120 success: `1.0`
- Route full483 probe success: `0.4741`
- Route full483 probe longest prefix: `170`

## Audit Conclusion

The final package can be generated from existing repo artifacts and documentation. The package must state clearly that full holder1-to-holder8 Gazebo transport is not solved.
"""


def official_artifacts_md() -> str:
    qwen = read_json("artifacts/v5/qwen_l1_demo/l1_to_rl_skill_request_qwen.json") or {}
    tool_call = qwen.get("tool_call", {"tool": "resolve_intent_packet", "arguments": {"object_id": "tray1", "source_slot": "shelf_A1", "target_slot": "shelf_B1", "constraints": {"speed_cap": "SLOW"}}})
    skill = qwen.get("skill_request", {})
    target_pose = skill.get("target_pose", {"xyz": [-0.92, -1.16, 1.22], "rpy": [3.14, 0.0, 3.14]})
    rows = "\n".join(
        f"| {s} | {succ:.2f} | {hp:.2f} mm | {ho:.4f} rad | {fp:.2f} mm | {fo:.4f} rad |"
        for s, succ, hp, ho, fp, fo in OFFICIAL["stage_rows"]
    )
    r = OFFICIAL["route"]
    return f"""
# Official Artifacts and Numbers

Purpose: This is the single source of truth for final report, slides, demo scripts, and README claims.

## A. Phase 1 Baseline

Phase 1 was operational but incomplete. The custom Gymnasium kinematic environment, training loop, throughput checks, and deterministic evaluation tools were running. Isolated approach and dock policies improved, but switched integration was unreliable and the VLM/Qwen layer had not started.

## B. Phase 2 Skill Stack

Final skill path:

```text
Approach -> Finisher
```

Dock-Coarse, Bridge, readiness classifier, acceptance map, and finisher adaptation are diagnostic steps. They helped identify the clean final path, but they are not the main final controller.

| Stage | Success | Handoff Pos Error | Handoff Ori Error | Final Pos Error | Final Ori Error |
|---:|---:|---:|---:|---:|---:|
{rows}

Core Phase 2 result:

```text
Stage 5 success rate: 0.93
Stage 5 handoff position error: 1.96 mm
Stage 5 handoff orientation error: 0.0177 rad
Stage 5 final position error: 2.89 mm
Stage 5 final orientation error: 0.0208 rad
```

## C. Qwen L1 Bridge

Demo command:

```text
{OFFICIAL["qwen_command"]}
```

Tool call:

```json
{json.dumps(tool_call, indent=2)}
```

Resolved skill request:

```text
object_id: {skill.get("object_id", "tray1")}
source_slot: {skill.get("source_slot", "shelf_A1")}
target_slot: {skill.get("target_slot", "shelf_B1")}
pipeline: {skill.get("pipeline", "APPROACH -> FINISHER")}
target xyz: {target_pose.get("xyz", [-0.92, -1.16, 1.22])}
target rpy: {target_pose.get("rpy", [3.14, 0.0, 3.14])}
```

Safety boundary:

```text
Qwen / L1 may produce semantic intent and structured skill requests.
Qwen / L1 must not produce raw joint actions, trajectories, torques, or delta_q.
L2/L3 own policy rollout, execution, and safety.
```

## D. Route Curriculum Result

Baseline:

```text
full483 baseline success rate: {r["baseline_success"]}
baseline longest prefix: {r["baseline_prefix"]}
```

Official route checkpoint:

```text
{r["best_checkpoint"]}
```

Best route result:

```text
prefix120 sequential success: {r["prefix120_success"]}
prefix120 longest prefix: {r["prefix120_prefix"]}
prefix120 route distance: {r["prefix120_distance"]:.3f} m
prefix120 mean final position error: {r["prefix120_pos_error"]:.5f} m
prefix120 mean final orientation error: {r["prefix120_ori_error"]:.5f} rad
full483 probe success: {r["full483_success"]}
full483 longest prefix: {r["full483_prefix"]}
full483 route distance: {r["full483_distance"]:.3f} m
first failure index: {r["first_failure_index"]}
first failure reason: {r["first_failure_reason"]}
```

## E. Failed Directions

These are important research findings, not final claims:

- Prefix180 direct fine-tune failed sequential retention.
- Prefix180 anti-forgetting reset-ratio retry failed sequential retention.
- Prefix180 teacher-anchor smoke was promising and reached prefix170, but the longer 1M latest checkpoint drifted.
- Dock-Coarse as a fixed middle layer was removed from the final main path.
- Full holder1-to-holder8 transport is not solved.
"""


def final_summary_md() -> str:
    return """
# Final Project Summary

Purpose: A concise landing page for the final ENPM690 project package.

## Project Title

Modular Three-Layer RL for Kitchen Manipulation

## One-Sentence Outcome

This project demonstrates a modular L1-to-RL manipulation prototype with a working kinematic RL skill stack, a Qwen semantic bridge, and a route-curriculum extension from local correction to long-prefix route following.

## What Is Completed

- A pure kinematic Gymnasium environment for a 7-DoF arm.
- PPO/TD3 training and deterministic evaluation infrastructure.
- A simplified and validated `Approach -> Finisher` skill stack.
- A Qwen MCP bridge that converts semantic commands into safe structured skill requests.
- A route curriculum system using dense q-goal targets as policy observations.
- A route-following improvement from longest prefix 21 to full-route probe prefix 170.

## What Is Not Completed

- Full holder1 -> holder8 transport.
- Full Gazebo physics validation.
- Real camera grounding.
- Contact, friction, and object interaction.
- Real robot deployment.

## Official Demos

- Demo 1: Qwen L1 bridge.
- Demo 2: kinematic `Approach -> Finisher` skill stack.
- Demo 3: route curriculum extension.

## Final Package Files

- `report/FINAL_REPORT.md`
- `report/FINAL_REPORT.pdf`
- `report/FINAL_PRESENTATION.pptx`
- `report/DEMO_VIDEO_SCRIPT.md`
- `report/OFFICIAL_ARTIFACTS.md`
"""


def final_report_md() -> str:
    rows = "\n".join(
        f"| {s} | {succ:.2f} | {hp:.2f} mm | {ho:.4f} rad | {fp:.2f} mm | {fo:.4f} rad |"
        for s, succ, hp, ho, fp, fo in OFFICIAL["stage_rows"]
    )
    return f"""
# {OFFICIAL["title"]}

**Course:** ENPM690 Final Project  
**Project:** Robot Brain Trainer  
**Author:** Che-Jung Chuang  
**Date:** May 2026

## 1. Abstract

This project implements a modular L1/L2/L3 robot learning stack for kitchen-style manipulation. L1 is a Qwen semantic bridge that turns a natural-language command into a validated `IntentPacket`; L2/L3 contain the learned kinematic manipulation stack and deterministic execution boundary. The Phase 2 controller was simplified to `Approach -> Finisher`, which reaches 93% success in the hardest current kinematic workspace stage with approximately 2.89 mm final position error and 0.0208 rad final orientation error. A route-curriculum extension further expands the learned controller from local waypoint correction to long-prefix route following: the original full-route probe had longest prefix 21, while the route-trained policy reaches prefix120 with 100% sequential success and reaches prefix170 in a full 483-waypoint probe. Full holder1-to-holder8 Gazebo kitchen transport is not yet solved.

## 2. Motivation

Kitchen manipulation is difficult to debug as a monolithic end-to-end system. Semantic interpretation, learned motor control, and low-level execution have different failure modes. This project therefore separates the stack into three layers:

- L1: semantic intent using Qwen / VLM-style reasoning.
- L2: learned RL skill policy.
- L3: deterministic execution and safety boundary.

This decomposition makes the system easier to inspect: the LLM cannot directly move the robot, and the motor policy does not need to solve natural-language understanding.

## 3. Original Proposal vs Final Scope

The original proposal targeted a modular kitchen manipulation pipeline:

```text
L1 semantic interpretation -> L2 learned skill -> L3 safe execution
```

The final scope is a structured proof-of-concept:

- A Qwen bridge that produces validated structured commands.
- A kinematic RL skill stack.
- Route-curriculum numeric validation.

The final system does not claim full Gazebo physics completion, real camera grounding, or real robot deployment.

## 4. System Architecture

![Final architecture](figures/final_architecture_l1_l2_l3.png)

The key contract is that L1 never emits raw controls. Qwen can choose object, source slot, target slot, and constraints. L2/L3 own policy inference, safety, and execution.

## 5. Phase 1 Baseline

Phase 1 established the research infrastructure:

- Gymnasium kinematic environment.
- PPO/TD3 training.
- Deterministic evaluation.
- Isolated approach and dock policies.

However, Phase 1 did not solve final integration. Approach/dock switching was unreliable, and the VLM/Qwen component had not started. Phase 1 proved infrastructure rather than final behavior.

## 6. Phase 2 Skill Stack

Phase 2 simplified the final control path to:

```text
Approach -> Finisher
```

Dock-Coarse, Bridge, readiness classifier, acceptance maps, and finisher adaptation were diagnostic tools. They helped reveal that adding modules was not always helpful; once Approach could produce clean handoff states, the cleanest final path was `Approach -> Finisher`.

| Stage | Success | Handoff Pos Error | Handoff Ori Error | Final Pos Error | Final Ori Error |
|---:|---:|---:|---:|---:|---:|
{rows}

The most demanding Stage 5 reached 93% success, with about 2.89 mm final position error and 0.0208 rad final orientation error.

![Workspace sweep](figures/workspace_sweep_stage_success.png)

## 7. Qwen L1 Bridge

The Qwen bridge exposes a structured tool interface:

- `get_l1_scene_context`
- `resolve_intent_packet`
- `prepare_phase1_skill_request`

For the command:

```text
{OFFICIAL["qwen_command"]}
```

Qwen produces a structured call to `resolve_intent_packet`, resolving:

```text
object_id: tray1
source_slot: shelf_A1
target_slot: shelf_B1
constraints: speed_cap = SLOW
resolved command: MOVE_PLATE(shelf_A1, shelf_B1)
pipeline: APPROACH -> FINISHER
target pose: [-0.92, -1.16, 1.22]
orientation: [3.14, 0.0, 3.14]
```

This validates the semantic-to-RL contract without letting the LLM emit low-level robot controls.

## 8. Route Curriculum Extension

The route curriculum was added after discovering that the policy was a strong local controller but not a full scene-level transport controller.

Before route curriculum:

```text
full483 success rate: 0.0435
longest continuous success prefix: 21
```

Implementation changes:

- Dense q-goal route target.
- Route-specific observation keys: `route_q_goal`, `route_q_error`, `route_tangent`, `route_scalar`.
- Prefix curriculum.
- Sequential actual-final-q evaluation.

After route curriculum:

```text
prefix120 sequential success: 1.0
prefix120 route distance: 1.720 m
full483 probe success: 0.4741
full483 longest prefix: 170
first failure index: 171
first failure reason: position
```

![Route prefix improvement](figures/route_prefix_improvement.png)

## 9. Negative Results and Lessons

Several negative results are central to the final interpretation:

- Prefix180 direct fine-tune failed early-prefix retention.
- Prefix180 anti-forgetting retry also failed.
- Sampled training success can be misleading.
- Sequential actual-final-q evaluation is the real route metric.
- Teacher-anchor smoke was promising, but a longer latest checkpoint drifted.

This means future training needs sequential-gated checkpoint selection, not only longer PPO fine-tuning.

![Route limitations](figures/route_curriculum_limitations.png)

## 10. Limitations

The system does not yet solve:

- Full holder1 -> holder8 route transport.
- Full Gazebo physics validation.
- Image-grounded Qwen-VL perception.
- Contact, friction, and tray dynamics.
- Real robot deployment.

## 11. Future Work

### Route Curriculum

- Sequential-gated checkpoint selection.
- Teacher-anchor early stopping.
- Prefix180 / prefix260 expansion.

### Gazebo Validation

- Move the route-trained checkpoint into Phase 3A runtime.
- Validate controller timing and joint execution.
- Add collision/contact/object stability checks.

### Perception

- Replace structured scene proxy with image-grounded Qwen-VL scene estimation.
- Project paper/object geometry into a ground-plane coordinate frame before using visual headings.

## 12. Conclusion

This project demonstrates a modular architecture, a working kinematic `Approach -> Finisher` skill stack, a Qwen semantic bridge, and a route-curriculum extension from local correction to long-prefix route following. Full kitchen transport remains future work, but the modular learning stack and route-curriculum direction are validated by quantitative results.
"""


def slide_source_md() -> str:
    return """
# Final Presentation Slide Source

1. Title: Modular L1-to-RL Manipulation Stack with Route-Curriculum Extension.
2. Problem Setup: kitchen manipulation needs semantic interpretation, learned control, and safe execution.
3. L1/L2/L3 Architecture: LLM/VLM never outputs raw joint commands.
4. Phase 1 Baseline: infrastructure worked; integrated behavior did not.
5. Phase 2 Simplification: diagnostics led to Approach -> Finisher.
6. Kinematic Skill Result: Stage 5 success 0.93, final position 2.89 mm, final orientation 0.0208 rad.
7. Qwen L1 Bridge: natural language to IntentPacket to dry-run skill request.
8. Route Curriculum Motivation: local policy worked; full route failed at prefix 21.
9. Route Curriculum Result: prefix120 success 1.0; full probe reaches prefix 170.
10. What Still Fails: full483 unsolved, prefix180 forgetting, need gated checkpoint selection.
11. Demo Structure: Qwen bridge, kinematic skill, route curriculum.
12. Conclusion: modular stack demonstrated; full Gazebo kitchen transport remains future work.
"""


def demo_video_script_md() -> str:
    return f"""
# Demo Video Script

Purpose: Script for recording the final ENPM690 demo.

## Demo 1: Qwen L1 Bridge

Command:

```text
{OFFICIAL["qwen_command"]}
```

Narration:

This demo shows the L1 semantic layer. Qwen reads a task command and produces a safe structured tool call. The MCP bridge validates the object, source, target, and constraints, then prepares an APPROACH -> FINISHER skill request. Notice that Qwen never outputs raw joint commands.

Run:

```bash
bash scripts/final/run_demo_01_qwen_bridge.sh
```

Generated video:

- `report/videos/demo_01_qwen_bridge.mp4`

## Demo 2: Kinematic Approach -> Finisher Skill

Narration:

This demo shows the learned Approach -> Finisher skill stack. Earlier versions required extra bridge or coarse docking modules, but ablation showed that the cleanest path is Approach -> Finisher. In the hardest current workspace stage, the policy reaches 93% success with millimeter-level final position error.

Run:

```bash
bash scripts/final/run_demo_02_kinematic_skill.sh
```

Generated video:

- `report/videos/demo_02_kinematic_skill.mp4`

## Demo 3: Route Curriculum

Narration:

This demo shows the route-curriculum extension. The original local controller failed after 21 dense route waypoints. After adding route-specific observations and prefix curriculum, the policy can complete the first 120 waypoints reliably and reaches 170 waypoints in a full-route probe. The full holder1-to-holder8 task remains unsolved, but the route curriculum clearly expands the learned controller's reliable route coverage.

Run:

```bash
bash scripts/final/run_demo_03_route_curriculum.sh
```

Generated video:

- `report/videos/demo_03_route_curriculum.mp4`

Full compiled video:

- `report/videos/final_demo_compilation.mp4`
"""


def demo_recording_md() -> str:
    return """
# Demo Recording Commands

Purpose: Practical recording checklist for the final demo.

## One-shot Dry Run

```bash
bash scripts/final/record_final_demo.sh
```

This creates:

- `report/demo_outputs/demo_01_qwen_bridge_output.json`
- `report/demo_outputs/demo_02_kinematic_skill_summary.json`
- `report/demo_outputs/demo_03_route_curriculum_summary.json`
- `report/demo_outputs/demo_run_summary.md`
- `report/videos/demo_01_qwen_bridge.mp4`
- `report/videos/demo_02_kinematic_skill.mp4`
- `report/videos/demo_03_route_curriculum.mp4`
- `report/videos/final_demo_compilation.mp4`

The MP4 files are generated headlessly from the official final-package outputs. They are suitable as report/demo artifacts even when Gazebo GUI recording is unavailable.

## Regenerate Videos Only

```bash
python3 scripts/final/generate_final_videos.py
```

## Manual Screen Recording

If OBS or the desktop recorder is available:

1. Open a terminal at the repo root.
2. Start recording.
3. Run `bash scripts/final/run_demo_01_qwen_bridge.sh`.
4. Show `report/demo_outputs/demo_01_qwen_bridge_output.json`.
5. Run `bash scripts/final/run_demo_02_kinematic_skill.sh`.
6. Show `report/figures/workspace_sweep_stage_success.png`.
7. Run `bash scripts/final/run_demo_03_route_curriculum.sh`.
8. Show `report/figures/route_prefix_improvement.png` and `report/figures/route_curriculum_limitations.png`.

## Notes

The scripts are designed to work without GUI or Gazebo. If `ffmpeg`/`imageio` video rendering is unavailable, record the terminal and figures manually using the commands above.
"""


def build_pdf(markdown_path: Path, pdf_path: Path) -> None:
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name="SmallMono", fontName="Courier", fontSize=8, leading=10, spaceAfter=5))
    styles["Heading1"].fontSize = 18
    styles["Heading2"].fontSize = 14
    story = []
    lines = markdown_path.read_text(encoding="utf-8").splitlines()
    in_code = False
    code_buf = []
    for line in lines:
        if line.startswith("```"):
            if not in_code:
                in_code = True
                code_buf = []
            else:
                in_code = False
                text = "<br/>".join(code_buf) or " "
                story.append(Paragraph(text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"), styles["SmallMono"]))
                story.append(Spacer(1, 0.08 * inch))
            continue
        if in_code:
            code_buf.append(line)
            continue
        if line.startswith("![") and "](" in line:
            img_rel = line.split("](", 1)[1].rstrip(")")
            img_path = markdown_path.parent / img_rel
            if img_path.exists():
                story.append(Image(str(img_path), width=6.2 * inch, height=3.2 * inch, kind="proportional"))
                story.append(Spacer(1, 0.12 * inch))
            continue
        if line.startswith("# "):
            story.append(Paragraph(line[2:], styles["Title"]))
            story.append(Spacer(1, 0.14 * inch))
        elif line.startswith("## "):
            story.append(Paragraph(line[3:], styles["Heading1"]))
        elif line.startswith("### "):
            story.append(Paragraph(line[4:], styles["Heading2"]))
        elif line.startswith("|"):
            # Keep markdown tables readable as monospaced text.
            story.append(Paragraph(line.replace("|", " | "), styles["SmallMono"]))
        elif line.startswith("- "):
            story.append(Paragraph("• " + line[2:], styles["BodyText"]))
        elif not line.strip():
            story.append(Spacer(1, 0.06 * inch))
        else:
            safe = line.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(safe, styles["BodyText"]))
    doc = SimpleDocTemplate(str(pdf_path), pagesize=letter, rightMargin=54, leftMargin=54, topMargin=48, bottomMargin=48)
    doc.build(story)


def add_textbox(slide, x, y, w, h, text, size=22, bold=False, color=(30, 30, 30)):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return box


def add_bullets(slide, x, y, w, h, bullets, size=18):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(size)
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.55, 0.35, 11.8, 0.55, title, 26, True, (38, 70, 83))
    if subtitle:
        add_textbox(slide, 0.58, 0.92, 11.2, 0.3, subtitle, 12, False, (80, 80, 80))


def build_pptx(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def slide_bg(slide, color=(248, 246, 238)):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*color)
        shape.line.fill.background()
        slide.shapes._spTree.remove(shape._element)
        slide.shapes._spTree.insert(2, shape._element)

    # 1
    s = prs.slides.add_slide(blank); slide_bg(s, (245, 240, 225))
    add_textbox(s, 0.7, 1.25, 11.5, 1.4, "Modular L1-to-RL\nManipulation Stack", 42, True, (38, 70, 83))
    add_textbox(s, 0.75, 3.0, 10.8, 0.7, "with Route-Curriculum Extension", 28, False, (231, 111, 81))
    add_textbox(s, 0.78, 5.65, 10.8, 0.6, "ENPM690 Final Project · Che-Jung Chuang · May 2026 · Robot Brain Trainer", 15, False)

    # 2
    s = prs.slides.add_slide(blank); slide_bg(s); add_title(s, "Problem Setup")
    add_bullets(s, 0.8, 1.35, 5.3, 3.5, ["Kitchen tasks combine semantics, motion, and execution safety.", "A plate/tray task needs object, source, target, pose, and constraints.", "A monolithic end-to-end loop hides failure causes."], 20)
    s.shapes.add_picture(str(FIG / "final_architecture_l1_l2_l3.png"), Inches(6.35), Inches(1.25), width=Inches(6.4))

    # 3
    s = prs.slides.add_slide(blank); slide_bg(s); add_title(s, "Original L1/L2/L3 Architecture")
    add_bullets(s, 0.8, 1.35, 11.6, 3.3, ["L1: Qwen semantic intent and structured tool calls.", "L2: learned RL skill policy.", "L3: deterministic execution and safety boundary.", "The LLM/VLM never outputs raw joint commands."], 22)

    # 4
    s = prs.slides.add_slide(blank); slide_bg(s); add_title(s, "Phase 1 Baseline")
    add_bullets(s, 0.8, 1.35, 11.6, 3.8, ["Gymnasium environment, training loop, and deterministic eval were operational.", "Approach and dock policies had isolated progress.", "Switched integration was unreliable.", "VLM/Qwen layer had not started."], 21)
    add_textbox(s, 0.9, 5.7, 10.8, 0.5, "Phase 1 proved infrastructure, not final integrated behavior.", 20, True, (231, 111, 81))

    # 5
    s = prs.slides.add_slide(blank); slide_bg(s); add_title(s, "Phase 2 Simplification")
    add_bullets(s, 0.8, 1.35, 11.6, 3.7, ["Explored Dock-Coarse, Bridge, classifier, acceptance map, and finisher adaptation.", "Ablation showed extra modules could degrade clean handoff states.", "Final mainline simplified to Approach -> Finisher."], 22)

    # 6
    s = prs.slides.add_slide(blank); slide_bg(s); add_title(s, "Kinematic Skill Result", "Hardest current workspace stage reaches 93% success.")
    s.shapes.add_picture(str(FIG / "workspace_sweep_stage_success.png"), Inches(0.75), Inches(1.25), width=Inches(6.1))
    add_bullets(s, 7.2, 1.6, 5.3, 3.2, ["Stage 5 success: 0.93", "Final position error: 2.89 mm", "Final orientation error: 0.0208 rad", "Skill stack works in trained kinematic workspace."], 21)

    # 7
    s = prs.slides.add_slide(blank); slide_bg(s); add_title(s, "Qwen L1 Bridge")
    add_bullets(s, 0.8, 1.25, 11.5, 4.0, ["Command: move tray1 from shelf_A1 to shelf_B1.", "Qwen calls resolve_intent_packet.", "Bridge returns IntentPacket and APPROACH -> FINISHER request.", "Qwen does not emit raw joint actions."], 21)

    # 8
    s = prs.slides.add_slide(blank); slide_bg(s); add_title(s, "Route Curriculum Motivation")
    add_bullets(s, 0.8, 1.35, 11.4, 3.7, ["The policy was precise locally but failed scene-level transport.", "Dense q_goal is used as target observation, not command.", "Original full route success: 0.0435.", "Original longest prefix: 21."], 22)

    # 9
    s = prs.slides.add_slide(blank); slide_bg(s); add_title(s, "Route Curriculum Result")
    s.shapes.add_picture(str(FIG / "route_prefix_improvement.png"), Inches(0.7), Inches(1.25), width=Inches(6.4))
    add_bullets(s, 7.35, 1.6, 5.2, 3.2, ["Prefix120 sequential success: 1.0", "Full probe success: 0.4741", "Full probe longest prefix: 170", "Route curriculum expands reliable route coverage."], 21)

    # 10
    s = prs.slides.add_slide(blank); slide_bg(s); add_title(s, "What Still Fails")
    s.shapes.add_picture(str(FIG / "route_curriculum_limitations.png"), Inches(0.7), Inches(1.25), width=Inches(6.4))
    add_bullets(s, 7.35, 1.55, 5.2, 3.6, ["Full 483-waypoint route is not solved.", "Prefix180 fine-tuning can forget early route behavior.", "Sampled success can be misleading.", "Need sequential-gated checkpoint selection."], 20)

    # 11
    s = prs.slides.add_slide(blank); slide_bg(s); add_title(s, "Demo Structure")
    add_bullets(s, 0.8, 1.35, 11.6, 3.8, ["Demo 1: Qwen bridge -> structured IntentPacket.", "Demo 2: kinematic Approach -> Finisher skill stack.", "Demo 3: route curriculum numeric validation.", "Each demo has a runnable script and fallback summary."], 22)

    # 12
    s = prs.slides.add_slide(blank); slide_bg(s); add_title(s, "Conclusion")
    add_bullets(s, 0.8, 1.35, 11.6, 4.0, ["Modular L1-to-RL stack demonstrated.", "Approach -> Finisher works in trained kinematic workspace.", "Qwen bridge validates semantic-to-RL interface.", "Route curriculum extends local correction to long-prefix route following.", "Full Gazebo kitchen transport remains future work."], 21)

    prs.save(path)


def demo_scripts() -> dict[str, str]:
    return {
        "run_demo_01_qwen_bridge.sh": r"""#!/usr/bin/env bash
set -euo pipefail
cd "$(dirname "$0")/../.."
mkdir -p report/demo_outputs
python3 - <<'PY'
import json
from pathlib import Path
src = Path("artifacts/v5/qwen_l1_demo/l1_to_rl_skill_request_qwen.json")
if not src.exists():
    print("WARN: Qwen artifact missing; using fallback summary.")
    payload = {"tool_call":{"tool":"resolve_intent_packet","arguments":{"object_id":"tray1","source_slot":"shelf_A1","target_slot":"shelf_B1","constraints":{"speed_cap":"SLOW"}}},"skill_request":{"pipeline":"APPROACH -> FINISHER","target_pose":{"xyz":[-0.92,-1.16,1.22],"rpy":[3.14,0.0,3.14]}}}
else:
    payload = json.loads(src.read_text())
skill = payload.get("skill_request", {})
summary = {
    "demo": "Qwen L1 bridge",
    "tool": payload.get("tool_call", {}).get("tool"),
    "object_id": skill.get("object_id", "tray1"),
    "source_slot": skill.get("source_slot", "shelf_A1"),
    "target_slot": skill.get("target_slot", "shelf_B1"),
    "pipeline": skill.get("pipeline", "APPROACH -> FINISHER"),
    "target_pose": skill.get("target_pose", {"xyz":[-0.92,-1.16,1.22],"rpy":[3.14,0.0,3.14]}),
    "safety_boundary": "L1 semantic only; no raw joint actions.",
}
out = Path("report/demo_outputs/demo_01_qwen_bridge_output.json")
out.write_text(json.dumps(summary, indent=2))
print(json.dumps(summary, indent=2))
PY
""",
        "run_demo_02_kinematic_skill.sh": r"""#!/usr/bin/env bash
set -euo pipefail
cd "$(dirname "$0")/../.."
mkdir -p report/demo_outputs report/figures
python3 scripts/final/generate_final_figures.py >/dev/null
python3 - <<'PY'
import json
from pathlib import Path
src = Path("artifacts/kinematic_phase1/phase1c/workspace_sweep_workspace_noop_vs_previous_summary_001.json")
summary = {"demo": "Kinematic Approach -> Finisher", "stage5_success": 0.93, "stage5_final_position_error_mm": 2.89, "stage5_final_orientation_error_rad": 0.0208}
if src.exists():
    payload = json.loads(src.read_text())
    row5 = next((r for r in payload.get("rows", []) if r.get("stage_index") == 5), None)
    if row5:
        new = row5.get("new", {})
        summary.update({
            "stage5_success": new.get("success_rate"),
            "stage5_handoff_position_error_mm": new.get("mean_handoff_position_error", 0) * 1000,
            "stage5_handoff_orientation_error_rad": new.get("mean_handoff_orientation_error"),
            "stage5_final_position_error_mm": new.get("mean_final_position_error", 0) * 1000,
            "stage5_final_orientation_error_rad": new.get("mean_final_orientation_error"),
        })
Path("report/demo_outputs/demo_02_kinematic_skill_summary.json").write_text(json.dumps(summary, indent=2))
print(json.dumps(summary, indent=2))
PY
""",
        "run_demo_03_route_curriculum.sh": r"""#!/usr/bin/env bash
set -euo pipefail
cd "$(dirname "$0")/../.."
mkdir -p report/demo_outputs report/figures
python3 scripts/final/generate_final_figures.py >/dev/null
python3 - <<'PY'
import json
from pathlib import Path
prefix = Path("artifacts/kinematic_phase1/route_curriculum/route_prefix120_routeobs_sequence2_1m_001/route_eval_sequential/route_eval_sequential_summary.json")
full = Path("artifacts/kinematic_phase1/route_curriculum/eval_prefix120_model_full483_001/route_eval_sequential_summary.json")
summary = {"demo": "Route curriculum", "baseline_success": 0.0435, "baseline_longest_prefix": 21, "prefix120_success": 1.0, "full483_success": 0.4741, "full483_longest_prefix": 170}
if prefix.exists():
    p = json.loads(prefix.read_text())
    summary.update({"prefix120_success": p.get("success_rate"), "prefix120_longest_prefix": p.get("longest_success_prefix"), "prefix120_distance_m": p.get("cumulative_successful_route_distance_m")})
if full.exists():
    f = json.loads(full.read_text())
    summary.update({"full483_success": f.get("success_rate"), "full483_longest_prefix": f.get("longest_success_prefix"), "first_failure_index": f.get("first_failure_index"), "first_failure_reason": f.get("first_failure_reason")})
Path("report/demo_outputs/demo_03_route_curriculum_summary.json").write_text(json.dumps(summary, indent=2))
print(json.dumps(summary, indent=2))
PY
""",
        "record_final_demo.sh": r"""#!/usr/bin/env bash
set -euo pipefail
cd "$(dirname "$0")/../.."
mkdir -p report/videos report/demo_outputs
bash scripts/final/run_demo_01_qwen_bridge.sh
bash scripts/final/run_demo_02_kinematic_skill.sh
bash scripts/final/run_demo_03_route_curriculum.sh
python3 scripts/final/generate_final_videos.py
cat > report/demo_outputs/demo_run_summary.md <<'EOF'
# Final Demo Run Summary

Generated demo summaries:

- `report/demo_outputs/demo_01_qwen_bridge_output.json`
- `report/demo_outputs/demo_02_kinematic_skill_summary.json`
- `report/demo_outputs/demo_03_route_curriculum_summary.json`

Generated demo videos:

- `report/videos/demo_01_qwen_bridge.mp4`
- `report/videos/demo_02_kinematic_skill.mp4`
- `report/videos/demo_03_route_curriculum.mp4`
- `report/videos/final_demo_compilation.mp4`

These are headless MP4 videos generated from the official final-package outputs.
EOF
echo "Demo summaries generated in report/demo_outputs"
echo "Demo videos generated in report/videos"
""",
        "check_final_package.sh": r"""#!/usr/bin/env bash
set -euo pipefail
cd "$(dirname "$0")/../.."
status=0
check() {
  if [ -e "$1" ]; then echo "PASS $1"; else echo "FAIL $1"; status=1; fi
}
warn() {
  if [ -e "$1" ]; then echo "PASS $1"; else echo "WARN $1"; fi
}
check report/FINAL_PROJECT_SUMMARY.md
check report/OFFICIAL_ARTIFACTS.md
check report/FINAL_REPORT.md
check report/FINAL_REPORT.pdf
check report/FINAL_PRESENTATION.pptx
check report/DEMO_VIDEO_SCRIPT.md
check report/DEMO_RECORDING_COMMANDS.md
check report/figures/final_architecture_l1_l2_l3.png
check report/figures/workspace_sweep_stage_success.png
check report/figures/route_prefix_improvement.png
check report/figures/route_curriculum_limitations.png
check scripts/final/run_demo_01_qwen_bridge.sh
check scripts/final/run_demo_02_kinematic_skill.sh
check scripts/final/run_demo_03_route_curriculum.sh
check scripts/final/record_final_demo.sh
check scripts/final/generate_final_videos.py
check report/videos/demo_01_qwen_bridge.mp4
check report/videos/demo_02_kinematic_skill.mp4
check report/videos/demo_03_route_curriculum.mp4
check report/videos/final_demo_compilation.mp4
warn artifacts/kinematic_phase1/phase1c/workspace_sweep_workspace_noop_vs_previous_summary_001.json
warn artifacts/v5/qwen_l1_demo/l1_to_rl_skill_request_qwen.json
warn artifacts/kinematic_phase1/route_curriculum/route_prefix120_routeobs_sequence2_1m_001/model_latest.zip
if grep -q "FINAL_PROJECT_SUMMARY.md" README.md; then echo "PASS README final package links"; else echo "FAIL README final package links"; status=1; fi
if [ "$status" -eq 0 ]; then echo "FINAL PACKAGE CHECK: PASS"; else echo "FINAL PACKAGE CHECK: FAIL"; fi
exit "$status"
""",
    }


def update_readme() -> None:
    path = ROOT / "README.md"
    text = path.read_text(encoding="utf-8") if path.exists() else "# Robot Brain Trainer\n"
    marker = "## Final Project Status"
    section = """
## Final Project Status

This project demonstrates a modular L1-to-RL manipulation prototype. The final kinematic skill path is `Approach -> Finisher`; the Qwen MCP bridge produces safe structured skill requests; and route curriculum extends longest route prefix coverage from 21 to 170 in a full-route probe. Full holder1 -> holder8 Gazebo transport is not yet solved.

Final package:

- [Final Project Summary](report/FINAL_PROJECT_SUMMARY.md)
- [Final Report PDF](report/FINAL_REPORT.pdf)
- [Final Presentation](report/FINAL_PRESENTATION.pptx)
- [Demo Video Script](report/DEMO_VIDEO_SCRIPT.md)
- [Official Artifacts](report/OFFICIAL_ARTIFACTS.md)
"""
    if marker in text:
        before = text.split(marker, 1)[0].rstrip()
        text = before + "\n\n" + section.strip() + "\n"
    else:
        text = text.rstrip() + "\n\n" + section.strip() + "\n"
    path.write_text(text, encoding="utf-8")


def main() -> None:
    ensure_dirs()
    subprocess.run(["python3", str(SCRIPTS / "generate_final_figures.py")], check=True, cwd=ROOT)
    write(REPORT / "FINAL_PACKAGE_AUDIT.md", audit_md())
    write(REPORT / "OFFICIAL_ARTIFACTS.md", official_artifacts_md())
    write(REPORT / "FINAL_PROJECT_SUMMARY.md", final_summary_md())
    write(REPORT / "FINAL_REPORT.md", final_report_md())
    write(REPORT / "FINAL_PRESENTATION.md", slide_source_md())
    write(REPORT / "DEMO_VIDEO_SCRIPT.md", demo_video_script_md())
    write(REPORT / "DEMO_RECORDING_COMMANDS.md", demo_recording_md())
    build_pdf(REPORT / "FINAL_REPORT.md", REPORT / "FINAL_REPORT.pdf")
    build_pptx(REPORT / "FINAL_PRESENTATION.pptx")
    for name, text in demo_scripts().items():
        path = SCRIPTS / name
        write(path, text)
        os.chmod(path, 0o755)
    update_readme()
    print("Final package generated under report/")


if __name__ == "__main__":
    main()
